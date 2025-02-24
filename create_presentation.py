from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import os
import cairosvg
from PIL import Image

def svg_to_png(svg_path, width=800, height=600):
    """
    Convert SVG to PNG using cairosvg
    """
    output_path = svg_path.replace('.svg', '.png')
    cairosvg.svg2png(url=svg_path, write_to=output_path, output_width=width, output_height=height)
    return output_path

def add_image_to_slide(slide, image_path, left, top, width, height):
    """
    Add an image to a slide with proper scaling
    """
    slide.shapes.add_picture(image_path, left, top, width, height)
    
    # Add caption
    txBox = slide.shapes.add_textbox(left, top + height + Inches(0.1), width, Inches(0.5))
    tf = txBox.text_frame
    tf.text = os.path.basename(image_path).replace('.png', '')
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_hallucinate_presentation(output_filename="hallucinate_presentation.pptx"):
    # First convert SVGs to PNGs
    datacenter_png = svg_to_png("datacenter_vs_edge.svg")
    architecture_png = svg_to_png("architecture.svg")

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # --- Slide 1: Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hallucinate: Edge-First MLOPS for Decentralized AI"
    subtitle.text = "A Framework for Low-Latency, Decentralized Machine Learning\nPowered by IPFS and libp2p"

    # --- Slide 2: The Challenge ---
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "The Challenge of Centralized AI"
    tf = body_shape.text_frame
    tf.text = "Current Limitations:"
    
    p = tf.add_paragraph()
    p.text = "• High latency due to datacenter-dependent compute"
    p = tf.add_paragraph()
    p.text = "• Network bandwidth bottlenecks"
    p = tf.add_paragraph()
    p.text = "• Scalability constraints with increasing demand"
    p = tf.add_paragraph()
    p.text = "• Dependency on continuous internet connectivity"
    p = tf.add_paragraph()
    p.text = "• Limited offline capabilities"

    # Add datacenter vs edge PNG
    add_image_to_slide(slide, datacenter_png, 
                      Inches(8), Inches(2), Inches(7), Inches(5))

    # --- Slide 3: Solution Overview ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = "Hallucinate: A New Paradigm"
    tf = body_shape.text_frame
    tf.text = "Key Features:"
    
    p = tf.add_paragraph()
    p.text = "• Edge-first architecture for minimal latency"
    p = tf.add_paragraph()
    p.text = "• Decentralized storage and computation"
    p = tf.add_paragraph()
    p.text = "• Offline-capable after initial model download"
    p = tf.add_paragraph()
    p.text = "• Seamless integration with existing ML frameworks"
    p = tf.add_paragraph()
    p.text = "• Resilient peer-to-peer communication"

    # --- Slide 4: Architecture Deep Dive ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = "Architecture Overview"
    tf = body_shape.text_frame
    tf.text = "Three-Layer Architecture:"
    
    # Add architecture PNG
    add_image_to_slide(slide, architecture_png,
                      Inches(8), Inches(2), Inches(7), Inches(5))

    # [Rest of the slides remain the same...]
    # --- Slide 5: Key Modules ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = "Key Modules"
    tf = body_shape.text_frame
    tf.text = "Core Components:"
    
    p = tf.add_paragraph()
    p.text = "IPFS_Accelerate_Py"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Model server endpoint multiplexer"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Supports CUDA, CPU, and OpenVINO inference"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "IPFS_Model_Manager_Py"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Intelligent model caching and retrieval"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Multi-source model downloads"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "IPFS_Kit_Py"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Local IPFS node management"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Tiered caching system"
    p.level = 1

    # --- Slide 6: Implementation Details ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = "Implementation Details"
    tf = body_shape.text_frame
    tf.text = "Technical Specifications:"
    
    p = tf.add_paragraph()
    p.text = "Model Support"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Hugging Face Transformers integration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• ONNX and OpenVINO optimization"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Network Layer"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• libp2p for peer discovery and routing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• IPFS for content-addressed storage"
    p.level = 1

    # --- Slide 7: Future Development ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = "Future Development"
    tf = body_shape.text_frame
    tf.text = "Roadmap:"
    
    p = tf.add_paragraph()
    p.text = "Near-term Goals"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• WebNN integration for browser-based inference"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Enhanced knowledge graph capabilities"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Long-term Vision"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Fully decentralized AI infrastructure"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Edge-first compute paradigm"
    p.level = 1

    # --- Slide 8: End Slide ---
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_slide_layout)

    # Add title
    left = Inches(1)
    top = Inches(2)
    width = Inches(14)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Thank You!\nPushing AI Compute to the Edge"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True

    # Add contact info
    left = Inches(2)
    top = Inches(4)
    width = Inches(12)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "GitHub: github.com/endomorphosis/hallucinate\nContact: hallucinate@example.com"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(24)

    # Save the presentation
    prs.save(output_filename)
    print(f"Presentation created successfully: {output_filename}")

    # Clean up temporary PNG files
    os.remove(datacenter_png)
    os.remove(architecture_png)

if __name__ == "__main__":
    create_hallucinate_presentation()
